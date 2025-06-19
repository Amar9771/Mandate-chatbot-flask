from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
title = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(8), Inches(0.8))
title_tf = title.text_frame
title_tf.text = "BWR Process Flowchart"
title_tf.paragraphs[0].font.size = Pt(32)
title_tf.paragraphs[0].font.bold = True

# Workflow steps
operational_steps = [
    "BD Brings Mandate",
    "Register Mandate & Upload Docs",
    "DDP – Financial Capture",
    "Allocation",
    "Validate the Data",
    "Update MIS & Upload Docs",
    "Peer Review:\n- Sector Reviewers\n- BWR Projections\n- Score Generation"
]

approval_steps = [
    "Prepare Rating Note / SH Approval",
    "QC (AI Optional)",
    "Committee Review",
    "Prepare Minutes / Rationale / SH Approval",
    "QC (2nd Level)",
    "Share Draft Rationale",
    "Publish"
]

# Layout config
left_x = Inches(1)
right_x = Inches(6)
top_y = Inches(1.2)
box_w = Inches(3.5)
box_h = Inches(0.9)
gap = Inches(0.3)
arrow_color = RGBColor(100, 100, 100)

def add_box(slide, x, y, w, h, text, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(255, 255, 255)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)
    return shape

# Add Operational Lane
left_boxes = []
for idx, text in enumerate(operational_steps):
    y = top_y + idx * (box_h + gap)
    box = add_box(slide, left_x, y, box_w, box_h, text, RGBColor(0, 102, 204))
    left_boxes.append(box)

# Add Approval Lane
right_boxes = []
for idx, text in enumerate(approval_steps):
    y = top_y + idx * (box_h + gap)
    box = add_box(slide, right_x, y, box_w, box_h, text, RGBColor(0, 153, 76))
    right_boxes.append(box)

# Add arrows between steps (within each lane)
for i in range(len(left_boxes) - 1):
    start = left_boxes[i]
    end = left_boxes[i + 1]
    slide.shapes.add_connector(
        connector_type=1,  # straight line
        begin_x=start.left + box_w / 2,
        begin_y=start.top + box_h,
        end_x=end.left + box_w / 2,
        end_y=end.top,
    ).line.color.rgb = arrow_color

for i in range(len(right_boxes) - 1):
    start = right_boxes[i]
    end = right_boxes[i + 1]
    slide.shapes.add_connector(
        connector_type=1,
        begin_x=start.left + box_w / 2,
        begin_y=start.top + box_h,
        end_x=end.left + box_w / 2,
        end_y=end.top,
    ).line.color.rgb = arrow_color

# Save it
output_path = r"C:\Users\Amareesh K\New\BWR_Professional_Flowchart.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
